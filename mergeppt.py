import sys, os, uuid, copy, random, subprocess, tempfile, threading
from pptx import Presentation
from pptx.dml.color import RGBColor
from lxml import etree
from PySide6.QtWidgets import (
    QApplication, QWidget, QVBoxLayout, QHBoxLayout,
    QListWidget, QListWidgetItem, QPushButton, QMessageBox, QFileDialog,
    QStyledItemDelegate, QStyle, QLabel, QLineEdit, QSplitter, QFrame,
    QColorDialog, QSpinBox, QCheckBox
)
from PySide6.QtCore import Qt, QSize, QRect, QRectF, QEvent, Signal, QUrl, QThread, QTimer
from PySide6.QtGui import (
    QPainter, QColor, QFont, QBrush, QPen, QFontMetrics, QPainterPath, QDrag,
    QDesktopServices
)
from PySide6.QtCore import QMimeData


# ── 병합 리스트 위젯 (내부 재정렬 + 외부 URL 드롭 지원) ──────────────────
class MergeListWidget(QListWidget):
    files_dropped = Signal(list)   # list of file path strings

    def __init__(self, parent=None):
        super().__init__(parent)
        self.setDragDropMode(QListWidget.DragDrop)
        self.setDefaultDropAction(Qt.MoveAction)
        self.setMouseTracking(True)
        self.viewport().setMouseTracking(True)

    def dragEnterEvent(self, event):
        if event.source() is self:
            event.setDropAction(Qt.MoveAction)
            event.accept()
        elif event.mimeData().hasUrls():
            event.setDropAction(Qt.CopyAction)
            event.accept()
        else:
            event.ignore()

    def dragMoveEvent(self, event):
        if event.source() is self:
            event.setDropAction(Qt.MoveAction)
            event.accept()
        elif event.mimeData().hasUrls():
            event.setDropAction(Qt.CopyAction)
            event.accept()
        else:
            event.ignore()

    def dropEvent(self, event):
        if event.source() is self:
            super().dropEvent(event)
        elif event.mimeData().hasUrls():
            drop_row = self._row_at(event.position().toPoint())
            paths = [
                url.toLocalFile() for url in event.mimeData().urls()
                if url.toLocalFile().lower().endswith(('.ppt', '.pptx'))
            ]
            if paths:
                self.files_dropped.emit((paths, drop_row))
            event.accept()
        else:
            event.ignore()

    def _row_at(self, pos):
        item = self.itemAt(pos)
        if item:
            return self.row(item)
        return self.count()


# ── 검색 결과 리스트 위젯 (드래그 시 URL 미임 데이터 제공) ─────────────────
class SearchResultsListWidget(QListWidget):
    def __init__(self, parent=None):
        super().__init__(parent)
        self.setDragEnabled(True)
        self.setDragDropMode(QListWidget.DragOnly)
        self.setSelectionMode(QListWidget.ExtendedSelection)

    def startDrag(self, supported_actions):
        items = self.selectedItems()
        if not items:
            return
        urls = [
            QUrl.fromLocalFile(it.data(Qt.UserRole))
            for it in items
            if it.data(Qt.UserRole)
        ]
        if not urls:
            return
        mime = QMimeData()
        mime.setUrls(urls)
        drag = QDrag(self)
        drag.setMimeData(mime)
        drag.exec(Qt.CopyAction)


# ── 백그라운드 파일 검색 스레드 ────────────────────────────────────────────
class FileSearchWorker(QThread):
    results_ready = Signal(list)   # list of (rel, full)

    def __init__(self, parent=None):
        super().__init__(parent)
        self._root  = ''
        self._query = ''
        self._stop_event = threading.Event()

    def start_search(self, root, query):
        self._stop_event.set()
        self.wait(400)
        self._root  = root
        self._query = query
        self._stop_event.clear()
        self.start()

    def stop(self):
        self._stop_event.set()

    def run(self):
        import unicodedata
        q = unicodedata.normalize('NFC', self._query).lower().strip()
        results = []
        try:
            for dirpath, dirnames, filenames in os.walk(self._root):
                if self._stop_event.is_set():
                    return
                dirnames[:] = [d for d in dirnames if not d.startswith('.')]
                for fn in filenames:
                    if self._stop_event.is_set():
                        return
                    if fn.lower().endswith(('.ppt', '.pptx')):
                        fn_nfc = unicodedata.normalize('NFC', fn).lower()
                        if not q or q in fn_nfc:
                            full = os.path.join(dirpath, fn)
                            rel  = os.path.relpath(full, self._root)
                            results.append((rel, full))
        except Exception:
            pass

        if self._stop_event.is_set():
            return

        self.results_ready.emit(sorted(results, key=lambda x: x[0].lower()))


# ── 아이템 델리게이트 (병합 리스트용) ─────────────────────────────────────
class PPTItemDelegate(QStyledItemDelegate):
    delete_requested = Signal(int)

    _H       = 52
    _M       = 3
    _HANDLE  = 24
    _NUM     = 28
    _BADGE   = 48
    _XBTN    = 32

    def paint(self, painter, option, index):
        painter.save()
        painter.setRenderHint(QPainter.Antialiasing)

        r = QRectF(option.rect.adjusted(self._M, self._M, -self._M, -self._M))
        selected = bool(option.state & QStyle.State_Selected)
        hovered  = bool(option.state & QStyle.State_MouseOver)

        if selected:
            bg, border = QColor("#2a3f5f"), QColor("#4a7abf")
        elif hovered:
            bg, border = QColor("#353550"), QColor("#4a4a66")
        else:
            bg, border = QColor("#2e2e42"), QColor("#3a3a52")

        path = QPainterPath()
        path.addRoundedRect(r, 8, 8)
        painter.fillPath(path, bg)
        painter.setPen(QPen(border, 1))
        painter.drawPath(path)

        file_path = index.data(Qt.UserRole) or ""
        name = os.path.basename(file_path)
        ext  = os.path.splitext(file_path)[1].lower()
        ri   = option.rect
        cy   = ri.center().y()

        # 햄버거 핸들
        f = QFont(); f.setPixelSize(18)
        painter.setFont(f)
        painter.setPen(QColor("#777" if hovered else "#555"))
        painter.drawText(
            QRect(ri.x() + 8, ri.y(), self._HANDLE, ri.height()),
            Qt.AlignCenter, "⠿"
        )

        # 번호
        f.setPixelSize(13)
        painter.setFont(f)
        painter.setPen(QColor("#888"))
        painter.drawText(
            QRect(ri.x() + 8 + self._HANDLE, ri.y(), self._NUM, ri.height()),
            Qt.AlignRight | Qt.AlignVCenter, f"{index.row() + 1}."
        )

        # X 버튼
        xr = self._xbtn_rect(ri)
        x_color = QColor("#ff6060") if hovered else QColor(180, 80, 80, 160)
        f.setPixelSize(13); f.setBold(True)
        painter.setFont(f)
        painter.setPen(x_color)
        painter.drawText(xr, Qt.AlignCenter, "✕")

        # 뱃지
        badge_right = xr.left() - 8
        badge_rect  = QRect(badge_right - self._BADGE, cy - 11, self._BADGE, 22)
        bp = QPainterPath()
        bp.addRoundedRect(QRectF(badge_rect), 4, 4)
        painter.fillPath(bp, QColor("#1a6bbf") if ext == '.pptx' else QColor("#c07000"))
        f.setPixelSize(11); f.setBold(True)
        painter.setFont(f)
        painter.setPen(QColor("white"))
        painter.drawText(badge_rect, Qt.AlignCenter, ext.lstrip('.').upper())

        # 파일명
        name_x = ri.x() + 8 + self._HANDLE + self._NUM + 8
        name_w = badge_rect.left() - name_x - 8
        f.setPixelSize(14); f.setBold(False)
        painter.setFont(f)
        painter.setPen(QColor("#e8e8f0"))
        elided = QFontMetrics(painter.font()).elidedText(name, Qt.ElideMiddle, name_w)
        painter.drawText(
            QRect(name_x, ri.y(), name_w, ri.height()),
            Qt.AlignLeft | Qt.AlignVCenter, elided
        )

        painter.restore()

    def _xbtn_rect(self, item_rect):
        r = item_rect
        return QRect(r.right() - self._XBTN - 4, r.y() + (r.height() - 26) // 2, 26, 26)

    def sizeHint(self, option, index):
        return QSize(0, self._H)

    def editorEvent(self, event, model, option, index):
        if event.type() == QEvent.Type.MouseButtonRelease:
            r = option.rect.adjusted(self._M, self._M, -self._M, -self._M)
            if self._xbtn_rect(r).contains(event.position().toPoint()):
                self.delete_requested.emit(index.row())
                return True
        return super().editorEvent(event, model, option, index)


# ── 메인 앱 ────────────────────────────────────────────────────────────────
class PPTMergerApp(QWidget):
    def __init__(self):
        super().__init__()
        self.bg_color    = QColor("#000000")  # 기본 배경색: 검정
        self.text_color  = QColor("#FFFFFF")  # 기본 글자색: 흰색
        self.slide_ratio = "16:9"             # 기본 슬라이드 비율
        self.text_valign = "top"              # 텍스트 세로 위치: top/center/bottom
        self.text_margin_pt = 20              # 텍스트 여백 (pt 단위)
        self._search_timer  = QTimer(self)
        self._search_timer.setSingleShot(True)
        self._search_timer.setInterval(350)
        self._search_timer.timeout.connect(self._do_search)
        self._search_worker = FileSearchWorker(self)
        self._search_worker.results_ready.connect(self._on_search_results)
        self.initUI()

    def initUI(self):
        self.setWindowTitle('PPT 병합기 ')
        self.resize(1080, 620)
        self.setAcceptDrops(True)
        self.setStyleSheet("background-color: #1e1e2e;")

        root = QVBoxLayout(self)
        root.setContentsMargins(12, 6, 12, 10)
        root.setSpacing(4)

        # 제목 + 힌트를 한 줄로
        header_row = QHBoxLayout()
        header_row.setSpacing(10)
        title = QLabel("PPT 병합기")
        title.setStyleSheet("color: #c0c0d8; font-size: 12px; font-weight: bold;")
        title.setFixedHeight(20)
        header_row.addWidget(title)
        hint = QLabel("드래그 앤 드롭  •  순서 드래그로 변경  •  우측 검색 결과 드래그하여 추가")
        hint.setStyleSheet("color: #444; font-size: 10px;")
        hint.setFixedHeight(20)
        header_row.addWidget(hint, 1)
        root.addLayout(header_row)

        # ── 좌우 분할 스플리터 ──
        splitter = QSplitter(Qt.Horizontal)
        splitter.setHandleWidth(8)
        splitter.setStyleSheet("""
            QSplitter::handle { background: #2e2e42; border-radius: 4px; }
            QSplitter::handle:hover { background: #5555aa; }
        """)

        # ── 좌측: 병합 리스트 ──
        left_container = QWidget()
        left_container.setStyleSheet("background: transparent;")
        left_layout = QVBoxLayout(left_container)
        left_layout.setContentsMargins(0, 0, 0, 0)
        left_layout.setSpacing(10)

        self.listWidget = MergeListWidget()
        self.listWidget.files_dropped.connect(self._on_files_dropped_to_list)
        self.listWidget.setStyleSheet("""
            QListWidget {
                background: #252535;
                border: 2px dashed #44445a;
                border-radius: 12px;
                padding: 4px;
                outline: none;
            }
            QListWidget::item { background: transparent; border: none; }
        """)
        self.delegate = PPTItemDelegate()
        self.delegate.delete_requested.connect(self._delete_row)
        self.listWidget.setItemDelegate(self.delegate)
        self.listWidget.model().rowsMoved.connect(
            lambda: self.listWidget.viewport().update()
        )
        left_layout.addWidget(self.listWidget)

        # ── 슬라이드 설정 행 (배경색 + 비율) ──
        settings_row = QHBoxLayout()
        settings_row.setSpacing(8)

        bg_lbl = QLabel("배경색:")
        bg_lbl.setStyleSheet("color: #888; font-size: 11px;")
        settings_row.addWidget(bg_lbl)

        self.bgColorBtn = QPushButton()
        self.bgColorBtn.setFixedSize(40, 28)
        self.bgColorBtn.setCursor(Qt.PointingHandCursor)
        self.bgColorBtn.setToolTip("클릭하여 배경색 선택")
        self.bgColorBtn.clicked.connect(self._pick_bg_color)
        self._update_bg_btn_style()
        settings_row.addWidget(self.bgColorBtn)

        settings_row.addSpacing(8)

        text_lbl = QLabel("글자색:")
        text_lbl.setStyleSheet("color: #888; font-size: 11px;")
        settings_row.addWidget(text_lbl)

        self.textColorBtn = QPushButton()
        self.textColorBtn.setFixedSize(40, 28)
        self.textColorBtn.setCursor(Qt.PointingHandCursor)
        self.textColorBtn.setToolTip("클릭하여 글자색 선택")
        self.textColorBtn.clicked.connect(self._pick_text_color)
        self._update_text_btn_style()
        settings_row.addWidget(self.textColorBtn)

        self.textColorChk = QCheckBox("일괄 적용")
        self.textColorChk.setChecked(False)  # 기본: 원본 색상 유지
        self.textColorChk.setToolTip("체크 시 선택한 글자색을 모든 텍스트에 적용\n미체크 시 원본 색상 유지")
        self.textColorChk.setStyleSheet("""
            QCheckBox { color: #888; font-size: 11px; spacing: 4px; }
            QCheckBox::indicator {
                width: 14px; height: 14px;
                border: 1px solid #44445a; border-radius: 3px;
                background: #1e1e2e;
            }
            QCheckBox::indicator:checked {
                background: #0063cc; border-color: #0063cc;
            }
            QCheckBox:hover { color: #ccc; }
        """)
        settings_row.addWidget(self.textColorChk)

        settings_row.addSpacing(16)

        ratio_lbl = QLabel("슬라이드 비율:")
        ratio_lbl.setStyleSheet("color: #888; font-size: 11px;")
        settings_row.addWidget(ratio_lbl)

        self.btn_ratio_169 = QPushButton("16:9")
        self.btn_ratio_169.setFixedSize(52, 28)
        self.btn_ratio_169.setCursor(Qt.PointingHandCursor)
        self.btn_ratio_169.clicked.connect(lambda: self._set_ratio("16:9"))
        settings_row.addWidget(self.btn_ratio_169)

        self.btn_ratio_43 = QPushButton("4:3")
        self.btn_ratio_43.setFixedSize(52, 28)
        self.btn_ratio_43.setCursor(Qt.PointingHandCursor)
        self.btn_ratio_43.clicked.connect(lambda: self._set_ratio("4:3"))
        settings_row.addWidget(self.btn_ratio_43)

        settings_row.addStretch(1)
        self._update_ratio_btn_styles()
        left_layout.addLayout(settings_row)

        # ── 텍스트 위치 설정 행 ──
        pos_row = QHBoxLayout()
        pos_row.setSpacing(8)

        pos_lbl = QLabel("텍스트 위치:")
        pos_lbl.setStyleSheet("color: #888; font-size: 11px;")
        pos_row.addWidget(pos_lbl)

        self.btn_valign_top = QPushButton("상단")
        self.btn_valign_top.setFixedSize(52, 28)
        self.btn_valign_top.setCursor(Qt.PointingHandCursor)
        self.btn_valign_top.clicked.connect(lambda: self._set_valign("top"))
        pos_row.addWidget(self.btn_valign_top)

        self.btn_valign_center = QPushButton("가운데")
        self.btn_valign_center.setFixedSize(60, 28)
        self.btn_valign_center.setCursor(Qt.PointingHandCursor)
        self.btn_valign_center.clicked.connect(lambda: self._set_valign("center"))
        pos_row.addWidget(self.btn_valign_center)

        self.btn_valign_bottom = QPushButton("하단")
        self.btn_valign_bottom.setFixedSize(52, 28)
        self.btn_valign_bottom.setCursor(Qt.PointingHandCursor)
        self.btn_valign_bottom.clicked.connect(lambda: self._set_valign("bottom"))
        pos_row.addWidget(self.btn_valign_bottom)

        pos_row.addSpacing(16)

        margin_lbl = QLabel("여백:")
        margin_lbl.setStyleSheet("color: #888; font-size: 11px;")
        pos_row.addWidget(margin_lbl)

        self.marginSpinBox = QSpinBox()
        self.marginSpinBox.setRange(0, 300)
        self.marginSpinBox.setValue(self.text_margin_pt)
        self.marginSpinBox.setSuffix(" pt")
        self.marginSpinBox.setFixedSize(76, 28)
        self.marginSpinBox.setStyleSheet("""
            QSpinBox {
                color: #e0e0f0; background: #1e1e2e;
                border: 1px solid #44445a; border-radius: 6px;
                padding: 2px 6px; font-size: 12px;
            }
            QSpinBox:focus { border-color: #5555aa; }
            QSpinBox::up-button, QSpinBox::down-button {
                width: 16px; background: #2e2e42; border: none;
            }
            QSpinBox::up-button:hover, QSpinBox::down-button:hover {
                background: #3a3a55;
            }
        """)
        self.marginSpinBox.valueChanged.connect(
            lambda v: setattr(self, 'text_margin_pt', v)
        )
        pos_row.addWidget(self.marginSpinBox)

        pos_row.addStretch(1)
        self._update_valign_btn_styles()
        left_layout.addLayout(pos_row)

        # 하단 버튼
        btn_row = QHBoxLayout()
        btn_row.setSpacing(12)

        self.clearButton = QPushButton("전체 초기화")
        self.clearButton.setFixedHeight(46)
        self.clearButton.setCursor(Qt.PointingHandCursor)
        self.clearButton.setStyleSheet("""
            QPushButton {
                color: #aaa; background: #2e2e42;
                border: 1px solid #44445a; border-radius: 10px; font-size: 14px;
            }
            QPushButton:hover { color: #fff; background: #3a3a55; border-color: #6666aa; }
            QPushButton:pressed { background: #2a2a45; }
        """)
        self.clearButton.clicked.connect(self.listWidget.clear)
        btn_row.addWidget(self.clearButton, 1)

        self.mergeButton = QPushButton("최종 파일로 합치기")
        self.mergeButton.setFixedHeight(46)
        self.mergeButton.setCursor(Qt.PointingHandCursor)
        self.mergeButton.setStyleSheet("""
            QPushButton {
                color: #fff;
                background: qlineargradient(x1:0,y1:0,x2:1,y2:0,stop:0 #0063cc,stop:1 #0096ff);
                border: none; border-radius: 10px;
                font-size: 16px; font-weight: bold;
            }
            QPushButton:hover {
                background: qlineargradient(x1:0,y1:0,x2:1,y2:0,stop:0 #0074ee,stop:1 #22aaff);
            }
            QPushButton:pressed {
                background: qlineargradient(x1:0,y1:0,x2:1,y2:0,stop:0 #0052aa,stop:1 #007acc);
            }
            QPushButton:disabled { background: #3a3a50; color: #666; }
        """)
        self.mergeButton.clicked.connect(self.merge_ppts)
        btn_row.addWidget(self.mergeButton, 3)

        left_layout.addLayout(btn_row)
        splitter.addWidget(left_container)

        # ── 우측: 파일 검색 패널 ──
        right_container = QWidget()
        right_container.setObjectName("searchPanel")
        right_container.setStyleSheet("""
            QWidget#searchPanel {
                background: #23233a;
                border: 1px solid #35355a;
                border-radius: 12px;
            }
        """)
        right_layout = QVBoxLayout(right_container)
        right_layout.setContentsMargins(14, 14, 14, 14)
        right_layout.setSpacing(8)

        # 검색 패널 제목
        search_title = QLabel("파일 검색")
        search_title.setStyleSheet(
            "color: #e0e0f0; font-size: 15px; font-weight: bold;"
            " background: transparent; border: none;"
        )
        right_layout.addWidget(search_title)

        sep = QFrame()
        sep.setFrameShape(QFrame.HLine)
        sep.setFixedHeight(1)
        sep.setStyleSheet("background: #35355a; border: none;")
        right_layout.addWidget(sep)

        # 폴더 경로
        lbl_folder = QLabel("검색 폴더")
        lbl_folder.setStyleSheet("color: #888; font-size: 11px; background: transparent; border: none;")
        right_layout.addWidget(lbl_folder)

        folder_row = QHBoxLayout()
        folder_row.setSpacing(6)
        self.folderInput = QLineEdit()
        self.folderInput.setPlaceholderText("최상위 폴더 경로...")
        self.folderInput.setStyleSheet(self._input_style())
        self.folderInput.returnPressed.connect(self._do_search)
        folder_row.addWidget(self.folderInput)

        browse_btn = QPushButton("…")
        browse_btn.setFixedSize(36, 32)
        browse_btn.setCursor(Qt.PointingHandCursor)
        browse_btn.setStyleSheet("""
            QPushButton {
                color: #bbb; background: #2e2e42;
                border: 1px solid #44445a; border-radius: 6px; font-size: 16px;
            }
            QPushButton:hover { color: #fff; background: #3a3a55; }
            QPushButton:pressed { background: #252540; }
        """)
        browse_btn.clicked.connect(self._browse_folder)
        folder_row.addWidget(browse_btn)
        right_layout.addLayout(folder_row)

        # 검색어 입력
        lbl_search = QLabel("파일명 검색")
        lbl_search.setStyleSheet("color: #888; font-size: 11px; background: transparent; border: none;")
        right_layout.addWidget(lbl_search)

        self.searchInput = QLineEdit()
        self.searchInput.setPlaceholderText("파일명을 입력하세요... (빈칸 = 전체)")
        self.searchInput.setStyleSheet(self._input_style())
        self.searchInput.textChanged.connect(self._on_search_changed)
        right_layout.addWidget(self.searchInput)

        # 결과 카운트
        self.resultCountLabel = QLabel("검색 결과: —")
        self.resultCountLabel.setStyleSheet(
            "color: #666; font-size: 11px; background: transparent; border: none;"
        )
        right_layout.addWidget(self.resultCountLabel)

        # 검색 결과 리스트
        self.searchList = SearchResultsListWidget()
        self.searchList.setStyleSheet("""
            QListWidget {
                background: #1e1e2e;
                border: 1px solid #35355a;
                border-radius: 8px;
                padding: 2px;
                outline: none;
            }
            QListWidget::item {
                color: #d0d0e8;
                padding: 5px 8px;
                border-bottom: 1px solid #252540;
                font-size: 12px;
                border-radius: 0px;
            }
            QListWidget::item:hover { background: #2a2a45; }
            QListWidget::item:selected { background: #2a3f5f; color: #ffffff; }
        """)
        self.searchList.itemDoubleClicked.connect(self._on_search_item_double_clicked)
        right_layout.addWidget(self.searchList, 1)

        drag_hint = QLabel("← 드래그하여 추가  |  더블클릭으로 추가")
        drag_hint.setAlignment(Qt.AlignCenter)
        drag_hint.setStyleSheet(
            "color: #4a4a6a; font-size: 10px; background: transparent; border: none;"
        )
        right_layout.addWidget(drag_hint)

        splitter.addWidget(right_container)
        splitter.setSizes([640, 420])
        splitter.setCollapsible(1, False)

        root.addWidget(splitter, 1)

    # ── 스타일 헬퍼 ──────────────────────────────────────────────────────
    @staticmethod
    def _input_style():
        return """
            QLineEdit {
                color: #e0e0f0;
                background: #1e1e2e;
                border: 1px solid #44445a;
                border-radius: 6px;
                padding: 5px 8px;
                font-size: 12px;
            }
            QLineEdit:focus { border-color: #5555aa; }
        """

    # ── 배경색 / 비율 UI 헬퍼 ─────────────────────────────────────────────
    def _update_bg_btn_style(self):
        hex_col = self.bg_color.name()  # e.g. "#ffffff"
        # 밝은 색이면 테두리를 회색으로, 어두운 색이면 그 색상 계열로
        border = "#44445a" if self.bg_color.lightness() > 128 else self.bg_color.lighter(150).name()
        self.bgColorBtn.setStyleSheet(
            f"QPushButton {{ background: {hex_col}; border: 2px solid {border};"
            f" border-radius: 6px; }}"
            f"QPushButton:hover {{ border-color: #5555aa; }}"
        )

    def _pick_bg_color(self):
        color = QColorDialog.getColor(self.bg_color, self, "배경색 선택")
        if color.isValid():
            self.bg_color = color
            self._update_bg_btn_style()

    def _update_text_btn_style(self):
        hex_col = self.text_color.name()
        # 배경은 흰색, 글자색 견본을 원형으로 표시
        self.textColorBtn.setStyleSheet(
            f"QPushButton {{ background: #2e2e42; border: 2px solid #44445a;"
            f" border-radius: 6px; }}"
            f"QPushButton::after {{ }}"  # placeholder
            f"QPushButton:hover {{ border-color: #5555aa; }}"
        )
        # 버튼 텍스트에 색상 견본(■) 표시
        self.textColorBtn.setText("■")
        self.textColorBtn.setStyleSheet(
            f"QPushButton {{ color: {hex_col}; background: #2e2e42;"
            f" border: 2px solid #44445a; border-radius: 6px;"
            f" font-size: 18px; }}"
            f"QPushButton:hover {{ border-color: #5555aa; }}"
        )

    def _pick_text_color(self):
        color = QColorDialog.getColor(self.text_color, self, "글자색 선택")
        if color.isValid():
            self.text_color = color
            self._update_text_btn_style()

    def _set_valign(self, valign: str):
        self.text_valign = valign
        self._update_valign_btn_styles()

    def _update_valign_btn_styles(self):
        active = (
            "QPushButton { color: #fff; background: #0063cc;"
            " border: none; border-radius: 6px; font-size: 12px; font-weight: bold; }"
        )
        inactive = (
            "QPushButton { color: #aaa; background: #2e2e42;"
            " border: 1px solid #44445a; border-radius: 6px; font-size: 12px; }"
            "QPushButton:hover { color: #fff; background: #3a3a55; }"
        )
        self.btn_valign_top.setStyleSheet(
            active if self.text_valign == "top" else inactive)
        self.btn_valign_center.setStyleSheet(
            active if self.text_valign == "center" else inactive)
        self.btn_valign_bottom.setStyleSheet(
            active if self.text_valign == "bottom" else inactive)

    def _set_ratio(self, ratio: str):
        self.slide_ratio = ratio
        self._update_ratio_btn_styles()

    def _update_ratio_btn_styles(self):
        active = (
            "QPushButton { color: #fff; background: #0063cc;"
            " border: none; border-radius: 6px; font-size: 12px; font-weight: bold; }"
        )
        inactive = (
            "QPushButton { color: #aaa; background: #2e2e42;"
            " border: 1px solid #44445a; border-radius: 6px; font-size: 12px; }"
            "QPushButton:hover { color: #fff; background: #3a3a55; }"
        )
        self.btn_ratio_169.setStyleSheet(active if self.slide_ratio == "16:9" else inactive)
        self.btn_ratio_43.setStyleSheet(active if self.slide_ratio == "4:3" else inactive)

    # ── 폴더 탐색 ─────────────────────────────────────────────────────────
    def _browse_folder(self):
        folder = QFileDialog.getExistingDirectory(self, "검색 폴더 선택")
        if folder:
            self.folderInput.setText(folder)
            self._do_search()

    # ── 검색 트리거 (디바운스) ────────────────────────────────────────────
    def _on_search_changed(self, _text):
        self._search_timer.stop()
        self._search_timer.start()

    def _do_search(self):
        root = self.folderInput.text().strip()
        if not root or not os.path.isdir(root):
            self.searchList.clear()
            self.resultCountLabel.setText("검색 결과: 폴더를 먼저 선택하세요")
            return
        self.resultCountLabel.setText("검색 중...")
        self._search_worker.start_search(root, self.searchInput.text().strip())

    def _on_search_results(self, results):
        """results: list of (rel, full)"""
        self.searchList.clear()
        for rel, full in results:
            item = QListWidgetItem(rel)
            item.setData(Qt.UserRole, full)
            item.setToolTip(full)
            self.searchList.addItem(item)
        n = len(results)
        self.resultCountLabel.setText("검색 결과: 없음" if not n else f"검색 결과: {n}개")

    def _on_search_item_double_clicked(self, item):
        path = item.data(Qt.UserRole)
        if path:
            self.add_file(path)

    # ── 드래그앤 드롭 (윈도우 레벨 — 리스트 위젯 바깥 영역용) ──────────────
    def dragEnterEvent(self, event):
        if event.mimeData().hasUrls():
            event.accept()
        else:
            event.ignore()

    def dropEvent(self, event):
        for url in event.mimeData().urls():
            path = url.toLocalFile()
            if path.lower().endswith(('.ppt', '.pptx')):
                self.add_file(path)

    # ── 파일 추가 / 삭제 ──────────────────────────────────────────────────
    def add_file(self, file_path):
        item = QListWidgetItem()
        item.setData(Qt.UserRole, file_path)
        self.listWidget.addItem(item)

    def _add_files_at(self, paths, row):
        for i, path in enumerate(paths):
            item = QListWidgetItem()
            item.setData(Qt.UserRole, path)
            self.listWidget.insertItem(row + i, item)

    def _on_files_dropped_to_list(self, payload):
        paths, row = payload
        self._add_files_at(paths, row)

    def _delete_row(self, row):
        self.listWidget.takeItem(row)

    # ── 병합 ──────────────────────────────────────────────────────────────
    def merge_ppts(self):
        if self.listWidget.count() == 0:
            QMessageBox.warning(self, "파일 없음", "병합할 PPT 파일을 추가해주세요.")
            return

        file_paths = [
            self.listWidget.item(i).data(Qt.UserRole)
            for i in range(self.listWidget.count())
        ]

        default_name = f"Merged_PPT_{uuid.uuid4().hex[:6]}.pptx"
        default_path = os.path.join(os.path.expanduser("~/Desktop"), default_name)

        merged_file_path, _ = QFileDialog.getSaveFileName(
            self, "저장 위치 선택", default_path, "PowerPoint 파일 (*.pptx)"
        )
        if not merged_file_path:
            return

        self.mergeButton.setEnabled(False)
        self.mergeButton.setText('변환 및 병합 중...')

        tmp_dir = tempfile.mkdtemp()
        try:
            converted = self._convert_ppt_files(file_paths, tmp_dir)

            merged = Presentation(converted[0])
            for src_path in converted[1:]:
                self._add_divider_slide(merged, self.bg_color)
                src_prs = Presentation(src_path)
                for slide in src_prs.slides:
                    self._copy_slide(merged, slide)

            # 슬라이드 크기 설정
            if self.slide_ratio == "16:9":
                merged.slide_width  = 12192000
                merged.slide_height = 6858000
            else:  # 4:3
                merged.slide_width  = 9144000
                merged.slide_height = 6858000

            slide_w = merged.slide_width
            slide_h = merged.slide_height

            # 슬라이드 마스터 / 레이아웃 이미지 제거 (근본 원인 차단)
            self._clean_slide_masters(merged)

            margin_emu = self.text_margin_pt * 12700  # pt → EMU

            # 전체 슬라이드 후처리
            for slide in merged.slides:
                self._remove_background_pictures(slide)
                self._fit_text_shapes(
                    slide, slide_w, slide_h, self.text_valign, margin_emu
                )

            # 배경색 적용
            self._apply_background_to_all_slides(merged, self.bg_color)
            # 글자색: 체크박스 선택 시에만 일괄 적용, 미선택 시 원본 유지
            if self.textColorChk.isChecked():
                self._apply_text_color_to_all_slides(merged, self.text_color)

            merged.save(merged_file_path)
            QMessageBox.information(self, "성공",
                f"병합 완료!\n'{merged_file_path}' 에 저장되었습니다.")
            QDesktopServices.openUrl(QUrl.fromLocalFile(merged_file_path))
        except Exception as e:
            QMessageBox.critical(self, "병합 실패", f"오류가 발생했습니다:\n{str(e)}")
        finally:
            import shutil
            shutil.rmtree(tmp_dir, ignore_errors=True)
            self.mergeButton.setEnabled(True)
            self.mergeButton.setText('최종 파일로 합치기')

    # ── .ppt → .pptx 변환 ────────────────────────────────────────────────
    @staticmethod
    def _get_soffice():
        import shutil
        if getattr(sys, 'frozen', False):
            base = sys._MEIPASS
            if sys.platform == 'darwin':
                return os.path.join(base, 'LibreOffice.app', 'Contents', 'MacOS', 'soffice')
            else:
                bundled = os.path.join(base, 'LibreOffice', 'program', 'soffice.exe')
                if os.path.isfile(bundled):
                    return bundled
                # bundled 경로에 없으면 시스템 설치 경로로 fallback
        if sys.platform == 'darwin':
            p = '/Applications/LibreOffice.app/Contents/MacOS/soffice'
            return p if os.path.exists(p) else (shutil.which('soffice') or p)
        if sys.platform == 'win32':
            candidates = [
                r'C:\Program Files\LibreOffice\program\soffice.exe',
                r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            ]
            try:
                import winreg
                for root in (winreg.HKEY_LOCAL_MACHINE, winreg.HKEY_CURRENT_USER):
                    for sub in (r'SOFTWARE\LibreOffice\LibreOffice',
                                r'SOFTWARE\WOW6432Node\LibreOffice\LibreOffice'):
                        try:
                            key = winreg.OpenKey(root, sub)
                            idx = 0
                            while True:
                                try:
                                    ver = winreg.EnumKey(key, idx)
                                    vkey = winreg.OpenKey(key, ver)
                                    path, _ = winreg.QueryValueEx(vkey, 'Path')
                                    candidates.append(os.path.join(path, 'program', 'soffice.exe'))
                                    idx += 1
                                except OSError:
                                    break
                        except OSError:
                            pass
            except Exception:
                pass
            for p in candidates:
                if os.path.exists(p):
                    return p
            found = shutil.which('soffice') or shutil.which('soffice.exe')
            if found:
                return found
            return None
        found = shutil.which('soffice')
        return found or 'soffice'

    def _convert_ppt_files(self, file_paths, tmp_dir):
        ppt_files = [p for p in file_paths if p.lower().endswith('.ppt')]
        if ppt_files:
            soffice = self._get_soffice()
            if not soffice or not os.path.isfile(soffice):
                raise FileNotFoundError(
                    "LibreOffice를 찾을 수 없습니다.\n"
                    "LibreOffice를 설치한 뒤 다시 시도해주세요.\n"
                    "(https://www.libreoffice.org/download/libreoffice-fresh/)\n\n"
                    "설치 후에도 이 오류가 발생하면 LibreOffice 설치 경로의 'program' 폴더가\n"
                    "시스템 PATH에 등록되어 있는지 확인해주세요."
                )
            kwargs = {}
            if sys.platform == 'win32':
                kwargs['creationflags'] = subprocess.CREATE_NO_WINDOW
            # LibreOffice 임시 사용자 프로필 지정 (번들 앱에서 쓰기 권한 문제 방지)
            lo_profile = os.path.join(tmp_dir, 'lo_profile')
            os.makedirs(lo_profile, exist_ok=True)
            cmd = [
                soffice,
                f'--env:UserInstallation=file://{lo_profile}',
                '--headless', '--norestore', '--nofirststartwizard',
                '--convert-to', 'pptx', '--outdir', tmp_dir,
            ] + ppt_files
            result_proc = subprocess.run(cmd, capture_output=True, text=True, **kwargs)
            if result_proc.returncode != 0:
                err_detail = (result_proc.stderr or result_proc.stdout or '').strip()
                raise RuntimeError(
                    f"LibreOffice 변환 실패 (종료 코드 {result_proc.returncode}).\n\n"
                    + (f"상세 오류:\n{err_detail}\n\n" if err_detail else "")
                    + "번들 앱에서 실행 시 터미널에서 아래 명령을 한 번 실행해 보세요:\n"
                    + f"xattr -cr \"{soffice}\""
                )
        result = []
        for p in file_paths:
            if p.lower().endswith('.ppt'):
                base = os.path.splitext(os.path.basename(p))[0]
                converted_path = os.path.join(tmp_dir, base + '.pptx')
                self._strip_slide_backgrounds(converted_path)
                result.append(converted_path)
            else:
                result.append(p)
        return result

    # 복사 시 destination이 자체 관리하는 관계 타입 — 건너뜀
    _SKIP_RELTYPES = {
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout',
        'http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide',
    }

    def _copy_slide(self, dest_prs, src_slide):
        blank = dest_prs.slide_layouts[min(6, len(dest_prs.slide_layouts) - 1)]
        dest_slide = dest_prs.slides.add_slide(blank)

        # slideLayout/notesSlide 제외한 모든 관계 복사
        rId_map = {}
        for rel_id, rel in src_slide.part.rels.items():
            if rel.reltype in self._SKIP_RELTYPES:
                continue
            try:
                if rel.is_external:
                    new_rid = dest_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
                else:
                    new_rid = dest_slide.part.relate_to(rel.target_part, rel.reltype)
                rId_map[rel_id] = new_rid
            except Exception:
                pass

        # 도형 트리 교체
        src_tree  = src_slide.shapes._spTree
        dest_tree = dest_slide.shapes._spTree
        for child in list(dest_tree):
            dest_tree.remove(child)
        for child in src_tree:
            dest_tree.append(copy.deepcopy(child))

        # r:id 매핑 반영
        if rId_map:
            xml = etree.tostring(dest_tree, encoding='unicode')
            for old, new in rId_map.items():
                xml = xml.replace(f'r:embed="{old}"', f'r:embed="{new}"')
                xml = xml.replace(f'r:id="{old}"',    f'r:id="{new}"')
                xml = xml.replace(f'r:link="{old}"',  f'r:link="{new}"')
            new_tree = etree.fromstring(xml)
            dest_tree.getparent().replace(dest_tree, new_tree)

        self._reassign_ids(dest_slide)

    @staticmethod
    def _strip_slide_backgrounds(pptx_path):
        """변환된 .ppt 파일의 각 슬라이드에서 명시적 배경(<p:bg>)을 제거해
        병합 후 첫 번째 파일의 테마를 상속받도록 한다."""
        ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            cSld = slide._element.find(f'{{{ns}}}cSld')
            if cSld is not None:
                bg = cSld.find(f'{{{ns}}}bg')
                if bg is not None:
                    cSld.remove(bg)
        prs.save(pptx_path)

    def _add_divider_slide(self, prs, color: QColor):
        """파일 사이 구분 슬라이드 삽입 — 선택한 배경색으로 채움."""
        slide = prs.slides.add_slide(prs.slide_layouts[min(6, len(prs.slide_layouts) - 1)])
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        sp_tree = slide.shapes._spTree
        keep = {f'{{{ns_p}}}nvGrpSpPr', f'{{{ns_p}}}grpSpPr'}
        for child in list(sp_tree):
            if child.tag not in keep:
                sp_tree.remove(child)

        hex_color = f"{color.red():02X}{color.green():02X}{color.blue():02X}"
        cSld = slide._element.find(f'{{{ns_p}}}cSld')
        existing = cSld.find(f'{{{ns_p}}}bg')
        if existing is not None:
            cSld.remove(existing)
        bg_xml = (
            f'<p:bg xmlns:p="{ns_p}" xmlns:a="{ns_a}">'
            f'<p:bgPr><a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'<a:effectLst/></p:bgPr></p:bg>'
        )
        cSld.insert(0, etree.fromstring(bg_xml))

    def _reassign_ids(self, slide):
        shape_id = 1
        for elem in slide._element.iter():
            local = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            if local == 'cNvPr':
                elem.set('id', str(shape_id))
                shape_id += 1
            if 'paraId' in elem.attrib:
                elem.set('paraId', format(random.randint(0x10000000, 0x7FFFFFFF), '08X'))
            if 'textId' in elem.attrib:
                elem.set('textId', format(random.randint(0x10000000, 0x7FFFFFFF), '08X'))

    def _apply_background_to_all_slides(self, prs, color: QColor):
        """모든 슬라이드에 단색 배경을 일괄 적용."""
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        hex_color = f"{color.red():02X}{color.green():02X}{color.blue():02X}"
        bg_xml = (
            f'<p:bg xmlns:p="{ns_p}" xmlns:a="{ns_a}">'
            f'<p:bgPr><a:solidFill><a:srgbClr val="{hex_color}"/></a:solidFill>'
            f'<a:effectLst/></p:bgPr></p:bg>'
        )
        for slide in prs.slides:
            cSld = slide._element.find(f'{{{ns_p}}}cSld')
            if cSld is not None:
                existing = cSld.find(f'{{{ns_p}}}bg')
                if existing is not None:
                    cSld.remove(existing)
                cSld.insert(0, etree.fromstring(bg_xml))

    @staticmethod
    def _clean_slide_masters(prs):
        """슬라이드 마스터 및 레이아웃의 spTree에서 이미지 도형을 모두 제거하고
        마스터 배경(p:bg)도 제거해 선택한 배경색만 표시되도록 한다."""
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        removable = {f'{{{ns_p}}}pic', f'{{{ns_p}}}sp', f'{{{ns_p}}}grpSp'}
        blip_tags = (f'{{{ns_a}}}blipFill', f'{{{ns_p}}}blipFill')
        text_tag  = f'{{{ns_a}}}t'

        def _strip_images(sp_tree):
            to_remove = [
                child for child in list(sp_tree)
                if child.tag in removable
                and any(child.find(f'.//{t}') is not None for t in blip_tags)
                and child.find(f'.//{text_tag}') is None
            ]
            for el in to_remove:
                sp_tree.remove(el)

        def _strip_bg(element):
            """cSld 안의 p:bg 요소 제거 (마스터/레이아웃 배경 이미지 차단)."""
            cSld = element.find(f'{{{ns_p}}}cSld')
            if cSld is not None:
                bg = cSld.find(f'{{{ns_p}}}bg')
                if bg is not None:
                    cSld.remove(bg)

        for master in prs.slide_masters:
            _strip_images(master.shapes._spTree)
            _strip_bg(master._element)
            for layout in master.slide_layouts:
                _strip_images(layout.shapes._spTree)
                _strip_bg(layout._element)

    @staticmethod
    def _remove_background_pictures(slide):
        """이미지 데이터(blipFill)가 있고 텍스트(a:t)가 없는 도형을
        배경 이미지로 판단해 제거. p:pic / p:sp / p:grpSp 모두 처리."""
        ns_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        ns_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'

        removable = {
            f'{{{ns_p}}}pic',
            f'{{{ns_p}}}sp',
            f'{{{ns_p}}}grpSp',
        }
        blip_tags = (f'{{{ns_a}}}blipFill', f'{{{ns_p}}}blipFill')
        text_tag  = f'{{{ns_a}}}t'

        sp_tree   = slide.shapes._spTree
        to_remove = []
        for child in list(sp_tree):
            if child.tag not in removable:
                continue
            has_image = any(child.find(f'.//{t}') is not None for t in blip_tags)
            has_text  = child.find(f'.//{text_tag}') is not None
            if has_image and not has_text:
                to_remove.append(child)

        for el in to_remove:
            sp_tree.remove(el)

    @staticmethod
    def _fit_text_shapes(slide, slide_width: int, slide_height: int,
                         valign: str = "center", margin: int = 0):
        """슬라이드에서 가장 큰 텍스트 박스(가사)만 위치·정렬을 적용하고,
        나머지 텍스트 박스는 경계 초과 시에만 보정한다.
        valign: 'top' | 'center' | 'bottom'
        margin: 슬라이드 경계와의 거리 (EMU)
        """
        text_shapes = [
            s for s in slide.shapes
            if s.has_text_frame
            and all(v is not None for v in (s.left, s.top, s.width, s.height))
        ]
        if not text_shapes:
            return

        # 가장 넓이가 큰 텍스트 박스 = 주요 가사 영역
        main = max(text_shapes, key=lambda s: int(s.width) * int(s.height))

        for shape in text_shapes:
            orig_w = int(shape.width)
            orig_h = int(shape.height)
            # 크기: 경계 초과 시에만 축소, 절대 키우지 않음
            w = min(orig_w, slide_width)
            h = min(orig_h, slide_height)

            if shape is main:
                # 가장 큰 박스만 좌우 가운데 + valign 적용
                l = (slide_width - w) // 2
                if valign == "top":
                    t = margin
                elif valign == "bottom":
                    t = slide_height - h - margin
                else:  # center
                    t = (slide_height - h) // 2
                t = max(0, min(t, slide_height - h))
                shape.left = l
                shape.top  = t
                if w < orig_w:
                    shape.width  = w
                if h < orig_h:
                    shape.height = h
            # 나머지(제목 등)는 위치·크기 일절 건드리지 않음

    def _apply_text_color_to_all_slides(self, prs, color: QColor):
        """모든 슬라이드의 텍스트 런(run)에 글자색을 일괄 적용."""
        rgb = RGBColor(color.red(), color.green(), color.blue())
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.color.rgb = rgb


if __name__ == '__main__':
    app = QApplication(sys.argv)
    app.setStyle('Fusion')
    app.setStyleSheet("""
        QMessageBox {
            background-color: #2e2e42;
        }
        QMessageBox QLabel {
            color: #e8e8f0;
            font-size: 14px;
        }
        QMessageBox QPushButton {
            color: #e8e8f0;
            background-color: #3a3a55;
            border: 1px solid #55558a;
            border-radius: 6px;
            padding: 6px 18px;
            font-size: 13px;
            min-width: 72px;
        }
        QMessageBox QPushButton:hover {
            background-color: #4a4a70;
            border-color: #7777bb;
        }
        QMessageBox QPushButton:pressed {
            background-color: #252540;
        }
    """)
    ex = PPTMergerApp()
    ex.show()
    sys.exit(app.exec())
